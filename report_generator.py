import matplotlib.pyplot as plt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from io import BytesIO
from PIL import Image
from pptx.util import Inches

def generate_charts():
    # Example: Generating a sample matplotlib chart
    plt.figure(figsize=(8, 5))
    plt.plot([1, 2, 3, 4], [10, 20, 25, 30], label="Sample Data")
    plt.title("Sample Chart")
    plt.legend()
    charts = [plt]
    return charts

def generate_pdf(data, charts):
    """Generate a PDF report."""
    pdf_file = "static/report.pdf"
    c = canvas.Canvas(pdf_file, pagesize=letter)
    width, height = letter

    # Add Title
    c.setFont("Helvetica-Bold", 18)
    c.drawString(100, height - 50, "Real-Time Report")

    # Add Data
    c.setFont("Helvetica", 12)
    y_position = height - 100
    for key, value in data.items():
        c.drawString(50, y_position, f"{key}: {value}")
        y_position -= 20

    # Add Charts
    y_position -= 30
    c.drawString(50, y_position, "Charts:")
    chart_y = y_position - 20
    for chart in charts:
        buf = BytesIO()
        chart.savefig(buf, format='png')
        buf.seek(0)
        img = Image.open(buf)
        img_width, img_height = img.size
        aspect = img_height / img_width
        img_width = 400
        img_height = img_width * aspect
        img = img.resize((int(img_width), int(img_height)))
        img_path = "temp_chart.png"
        img.save(img_path)

        c.drawImage(img_path, 50, chart_y - img_height, width=img_width, height=img_height)
        chart_y -= img_height + 20
        buf.close()

    c.save()
    return pdf_file

def generate_ppt(data, charts):
    """Generate a PowerPoint report."""
    ppt_file = "static/report.pptx"
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Real-Time Report"

    # Data Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    content = slide.placeholders[1]
    content.text = "\n".join([f"{key}: {value}" for key, value in data.items()])

    # Add Charts Slide
    for chart in charts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        buf = BytesIO()
        chart.savefig(buf, format='png')
        buf.seek(0)
        img_path = "chart_image.png"
        with open(img_path, 'wb') as f:
            f.write(buf.getvalue())
        slide.shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(6))
        buf.close()

    prs.save(ppt_file)
    return ppt_file
